"""Per-format parsers behind one registry. Phase 3.

Every parser returns the same shape: an ordered list of ``Page(number, text)``.
Citations deep-link to a page, so page provenance must survive the whole pipeline —
a parser that returns one blob of text has broken citations for that format.

Uploaded documents are untrusted input. Parsing happens in the worker, never in a
request handler.
"""

import io
import re
import zipfile
from dataclasses import dataclass
from typing import BinaryIO

from app.core.config import settings
from app.db.models.enums import SourceFormat


class UnparseableDocument(Exception):
    """A document that cannot be ingested, with a message fit for its owner.

    Deterministic: the same bytes fail the same way tomorrow, so the ingest task
    maps this straight to terminal failure instead of retrying it three times
    first.
    """

# Below this many characters per page on average, the document is probably scanned
# images rather than text. Indexing it would produce a tutor that confidently knows
# nothing, so it is flagged for OCR instead.
MIN_CHARS_PER_PAGE = 80

# TXT and MD have no pages, so one is synthesised every N characters purely so a
# citation can point somewhere a reader can find.
SYNTHETIC_PAGE_CHARS = 3000


@dataclass(frozen=True, slots=True)
class Page:
    number: int
    text: str


SNIFF_BYTES = 8192


def sniff_format(stream: BinaryIO, filename: str | None = None) -> SourceFormat | None:
    """Identify the format from the content.

    Never from the filename extension or the client's Content-Type: both are
    attacker-controlled. The filename is used only to break the txt/md tie, where the
    two are byte-identical and the distinction does not affect parsing.

    Takes a seekable stream rather than bytes. A ZIP container keeps its central
    directory at the END of the file, so DOCX and PPTX cannot be told apart from a
    leading slice — but zipfile seeks rather than reads, so handing it the spooled
    upload identifies the format without pulling the document into memory.
    """
    stream.seek(0)
    head = stream.read(SNIFF_BYTES)
    stream.seek(0)

    if head.startswith(b"%PDF-"):
        return SourceFormat.PDF

    # DOCX and PPTX are both ZIP containers; the entry names tell them apart. Any
    # other archive is a candidate for "one book uploaded in parts" (D26), decided
    # by looking at what the members actually hold — their names are as
    # attacker-controlled as the outer filename.
    if head.startswith(b"PK\x03\x04"):
        try:
            with zipfile.ZipFile(stream) as archive:
                names = archive.namelist()
                if any(name.startswith("word/") for name in names):
                    return SourceFormat.DOCX
                if any(name.startswith("ppt/") for name in names):
                    return SourceFormat.PPTX
                if _zip_holds_documents(archive):
                    return SourceFormat.ZIP
        except (zipfile.BadZipFile, OSError):
            return None
        finally:
            stream.seek(0)
        return None

    # Anything that decodes as UTF-8 without control bytes is treated as text.
    if b"\x00" in head:
        return None
    try:
        head.decode("utf-8")
    except UnicodeDecodeError:
        return None
    if filename and filename.lower().endswith((".md", ".markdown")):
        return SourceFormat.MD
    return SourceFormat.TXT


# --- ZIP: one book uploaded in parts (DECISIONS.md D26) ----------------------
#
# NCERT-style downloads arrive as an archive of per-chapter PDFs. The archive
# stays the stored source object; its members are combined here, at parse time,
# in reading order, with page numbers continuing across the seams — so a citation
# still points at exactly one place in the combined book, and re-ingest needs
# nothing but the original upload.

# How many member heads the upload-time sniff inspects before giving up.
_SNIFF_ZIP_MEMBERS = 50

# Decompression happens in bounded reads so a lying size header cannot balloon.
_ZIP_READ_BLOCK = 1 << 20

# Finder metadata in zips made on a Mac; never part of the book.
_ZIP_JUNK_PREFIX = "__MACOSX/"

# Runs of digits sort numerically, so `ch2.pdf` reads before `ch10.pdf`. A plain
# lexicographic sort would interleave the chapters of any book with more than
# nine parts, which is most textbooks.
_DIGIT_RUNS = re.compile(r"(\d+)")

# Same job as services.books.title_from_filename, scoped to a member name. Not
# imported from there: services depend on rag, never the other way around.
_MEMBER_SEPARATORS = re.compile(r"[_\-\s]+")


@dataclass(frozen=True, slots=True)
class ZipPart:
    """One member's place in the combined book.

    ``outline`` is the part's own top-level PDF outline with pages already made
    absolute, empty for non-PDF parts and PDFs without one. Chapter detection
    prefers it to the filename when it names real structure.
    """

    title: str
    page_start: int
    page_end: int
    outline: tuple[tuple[str, int], ...]


def _natural_key(name: str) -> tuple[int | str, ...]:
    return tuple(int(tok) if tok.isdigit() else tok.lower() for tok in _DIGIT_RUNS.split(name))


def _member_basename(name: str) -> str:
    return name.replace("\\", "/").rsplit("/", 1)[-1]


def _part_title(name: str) -> str:
    stem = _member_basename(name)
    if "." in stem:
        stem = stem.rsplit(".", 1)[0]
    cleaned = _MEMBER_SEPARATORS.sub(" ", stem).strip().lstrip(".").strip()
    return cleaned[:200] or "Part"


def _zip_holds_documents(archive: zipfile.ZipFile) -> bool:
    """Whether a generic ZIP looks like a book uploaded in parts.

    A shallow look at the first bytes of a bounded number of members, run on the
    request path where a full parse would block the event loop. The parse-time
    walk is the authority; this only turns away archives of obvious junk —
    images, executables — before a row is written for them. An encrypted member
    counts as a document so the refusal the owner eventually reads names the
    password, not a wrong file type.
    """
    checked = 0
    for info in archive.infolist():
        if checked >= _SNIFF_ZIP_MEMBERS:
            break
        if info.is_dir():
            continue
        basename = _member_basename(info.filename)
        if info.filename.startswith(_ZIP_JUNK_PREFIX) or basename.startswith("."):
            continue
        if info.flag_bits & 0x1:
            return True
        checked += 1
        try:
            with archive.open(info) as handle:
                head = handle.read(16)
        except (zipfile.BadZipFile, OSError, RuntimeError):
            continue
        # A PDF, a nested container (a DOCX, or a zip parse time will refuse with
        # a reason), or something texty enough to be TXT/MD.
        if head.startswith((b"%PDF-", b"PK\x03\x04")) or (head and b"\x00" not in head):
            return True
    return False


def _document_members(archive: zipfile.ZipFile) -> list[zipfile.ZipInfo]:
    """The members that could be documents, in reading order."""
    members: list[zipfile.ZipInfo] = []
    for info in archive.infolist():
        if info.is_dir():
            continue
        basename = _member_basename(info.filename)
        if info.filename.startswith(_ZIP_JUNK_PREFIX) or basename.startswith("."):
            continue
        if info.flag_bits & 0x1:
            raise UnparseableDocument(
                "That ZIP is password-protected. Remove the password and upload it again."
            )
        members.append(info)
    if len(members) > settings.zip_max_members:
        raise UnparseableDocument(
            f"That ZIP holds {len(members)} files; the limit is {settings.zip_max_members}."
        )
    return sorted(members, key=lambda info: _natural_key(info.filename))


def _pdf_extent(data: bytes) -> tuple[int, list[tuple[str, int]]]:
    """Page count and top-level outline, without extracting any text."""
    import fitz

    with fitz.open(stream=data, filetype="pdf") as doc:
        outline = doc.get_toc() or []
        count = doc.page_count
    tops = [(title.strip(), page) for level, title, page in outline if level == 1 and page > 0]
    return count, tops


def _scan_zip(data: bytes, *, with_text: bool) -> tuple[list[Page], list[ZipPart]]:
    """Walk the archive once, combining members into one paged book.

    ``with_text=False`` is the cheap pass chapter detection uses: part boundaries
    and outlines without pulling every page's text a second time — for PDFs, the
    bulk of any real book, that is a page count and a table of contents.

    Unsupported members are skipped so a stray README does not sink the book, but
    a member that *should* parse and does not fails loudly with its name — a book
    silently missing chapter 7 is worse than one that failed.
    """
    try:
        archive = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile as exc:
        raise UnparseableDocument("That ZIP file is damaged and could not be opened.") from exc

    pages: list[Page] = []
    parts: list[ZipPart] = []
    offset = 0
    remaining = settings.zip_max_uncompressed_bytes

    with archive:
        for info in _document_members(archive):
            blocks: list[bytes] = []
            with archive.open(info) as handle:
                while block := handle.read(_ZIP_READ_BLOCK):
                    remaining -= len(block)
                    if remaining < 0:
                        raise UnparseableDocument(
                            f"That ZIP unpacks past {settings.zip_max_uncompressed_mb} MB, "
                            "which is more than a book."
                        )
                    blocks.append(block)
            member = b"".join(blocks)

            member_format = sniff_format(io.BytesIO(member), info.filename)
            if member_format is None or member_format is SourceFormat.ZIP:
                # Not a document (or a nested archive, which is not recursed
                # into — one level of "in parts" is the product, D26).
                continue

            basename = _member_basename(info.filename)
            try:
                outline: tuple[tuple[str, int], ...] = ()
                if member_format is SourceFormat.PDF:
                    count, tops = _pdf_extent(member)
                    outline = tuple((title, offset + page) for title, page in tops)
                    member_pages = PARSERS[member_format](member) if with_text else []
                else:
                    member_pages = PARSERS[member_format](member)
                    count = len(member_pages)
            except UnparseableDocument:
                raise
            except Exception as exc:
                raise UnparseableDocument(
                    f'One file inside that ZIP could not be read: "{basename}". '
                    "Remove or replace it and upload again."
                ) from exc

            if count == 0:
                continue
            if with_text:
                pages.extend(
                    Page(number=offset + page.number, text=page.text) for page in member_pages
                )
            parts.append(
                ZipPart(
                    title=_part_title(info.filename),
                    page_start=offset + 1,
                    page_end=offset + count,
                    outline=outline,
                )
            )
            offset += count

    if not parts:
        raise UnparseableDocument(
            "No readable documents were found inside that ZIP. It should contain "
            "PDF, DOCX, PPTX, TXT or Markdown files."
        )
    return pages, parts


def _parse_zip(data: bytes) -> list[Page]:
    return _scan_zip(data, with_text=True)[0]


def zip_member_names(data: bytes) -> list[str]:
    """The archive's candidate members, in reading order, for the ingest trace.

    Reads the central directory only — no decompression — because this exists to
    answer "did it find all eighteen chapters, and in what order" on a screen, not
    to feed the pipeline. Returns an empty list rather than raising: a trace that
    fails is not a reason to fail a book that parsed.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            return [_member_basename(info.filename) for info in _document_members(archive)]
    except (zipfile.BadZipFile, OSError, UnparseableDocument):
        return []


def zip_outline(data: bytes) -> list[ZipPart]:
    """Part boundaries for chapter detection, pages absolute in the combined book.

    Chapter detection runs after parse and re-walks the archive rather than
    smuggling state between the two steps — the same trade every other format
    makes by handing ``detect_chapters`` the raw bytes.
    """
    return _scan_zip(data, with_text=False)[1]


def _parse_pdf(data: bytes) -> list[Page]:
    import fitz  # PyMuPDF

    with fitz.open(stream=data, filetype="pdf") as doc:
        return [Page(number=i + 1, text=page.get_text()) for i, page in enumerate(doc)]


def _parse_docx(data: bytes) -> list[Page]:
    """DOCX has no page concept until it is laid out, so pages are synthesised from
    paragraph runs. The number is stable for a given file, which is what a citation
    needs."""
    from docx import Document

    document = Document(io.BytesIO(data))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]

    pages: list[Page] = []
    buffer: list[str] = []
    for paragraph in paragraphs:
        buffer.append(paragraph)
        if sum(len(p) for p in buffer) >= SYNTHETIC_PAGE_CHARS:
            pages.append(Page(number=len(pages) + 1, text="\n\n".join(buffer)))
            buffer = []
    if buffer:
        pages.append(Page(number=len(pages) + 1, text="\n\n".join(buffer)))
    return pages


def _parse_pptx(data: bytes) -> list[Page]:
    """One slide, one page — the natural unit of provenance here."""
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    pages = []
    for number, slide in enumerate(presentation.slides, start=1):
        parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        pages.append(Page(number=number, text="\n".join(parts)))
    return pages


def _parse_text(data: bytes) -> list[Page]:
    content = data.decode("utf-8", errors="replace")
    return [
        Page(number=i + 1, text=content[start : start + SYNTHETIC_PAGE_CHARS])
        for i, start in enumerate(range(0, max(len(content), 1), SYNTHETIC_PAGE_CHARS))
    ]


PARSERS = {
    SourceFormat.PDF: _parse_pdf,
    SourceFormat.DOCX: _parse_docx,
    SourceFormat.PPTX: _parse_pptx,
    SourceFormat.TXT: _parse_text,
    SourceFormat.MD: _parse_text,
    SourceFormat.ZIP: _parse_zip,
}


def parse(data: bytes, source_format: SourceFormat) -> list[Page]:
    return PARSERS[source_format](data)


def looks_scanned(pages: list[Page]) -> bool:
    """True when there is too little text to be worth indexing."""
    if not pages:
        return True
    total = sum(len(page.text.strip()) for page in pages)
    return (total / len(pages)) < MIN_CHARS_PER_PAGE
